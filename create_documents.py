import os
import json
from pptx import Presentation
from pptx.util import Inches
from text_retrieval import create_vector_db
from image_retrieval import convert_pdfs_to_images, load_existing_image_mappings
from Chatbot.Mistral_7b import retrieve_faiss, retrieve_context
from transformers import pipeline, AutoProcessor
from byaldi import RAGMultiModalModel
from bedrock_handler import call_claude
from docx import Document
from typing import Dict, List


# ==================================================
# CONFIGURATION & PATHS
# ==================================================
DATA_PATH = "path to a folder called pdfs"
FAISS_DB_PATH = "path to a folder called vector_data_base"
IMAGES_FOLDER = "path to a folder called img_output"
OUTPUT_FOLDER = "path to a folder called Final_slide"

# ==================================================
# 1. PDF Conversion Caching
# ==================================================
def convert_pdfs_if_needed(data_path, images_folder):
    if not os.path.exists(images_folder) or not os.listdir(images_folder):
        all_images, file_names = convert_pdfs_to_images(data_path, images_folder)
        print("PDFs converted to images.")
    else:
        all_images = load_existing_image_mappings(images_folder)
        file_names = list(all_images.keys())
        print("Using cached image files.")
    return all_images, file_names

# ==================================================
# 2. Document Indexing Caching
# ==================================================
def index_documents_if_needed(model, data_path, index_name, force_reindex=False):
    cache_dir = ".cache"
    os.makedirs(cache_dir, exist_ok=True)
    index_indicator = os.path.join(cache_dir, f"{index_name}_index_indicator.txt")

    if force_reindex or not os.path.exists(index_indicator):
        model.index(input_path=data_path, index_name=index_name, overwrite=True)
        with open(index_indicator, "w") as f:
            f.write("indexed")
        print("Documents indexed.")
    else:
        print("Using cached document index.")
    return model

# ==================================================
# 3. Vector Database Creation/Loading Caching
# ==================================================
def create_or_load_vector_db(data_path, faiss_db_path, force_rebuild=False):
    if force_rebuild or not os.path.exists(faiss_db_path) or not os.listdir(faiss_db_path):
        os.makedirs(faiss_db_path, exist_ok=True)
        create_vector_db(data_path, faiss_db_path)
        print("Vector database created.")
    else:
        print("Using existing vector database.")

    db = retrieve_faiss(faiss_db_path)

    if hasattr(db, "index") and hasattr(db.index, "ntotal"):
        print(f"[INFO] FAISS DB contains {db.index.ntotal} vectors.")
        if db.index.ntotal == 0:
            raise ValueError("FAISS index is empty. Re-index with non-empty PDFs.")

    return db

# ==================================================
# 4. Model and Pipeline Initialization (Without Caching)
# ==================================================
def initialize_models():
    docs_retrieval_model = RAGMultiModalModel.from_pretrained("vidore/colpali-v1.2")
    model_id = "llava-hf/llava-1.5-7b-hf"
    pipe = pipeline("image-to-text", model=model_id, device=0)
    processor = AutoProcessor.from_pretrained(model_id)
    return docs_retrieval_model, pipe, processor

# ==================================================
# 5. Text Context Retrieval
# ==================================================
def retrieve_text_context(slide_headings_text, db, docs_retrieval_model):
    text_context = retrieve_context(slide_headings_text, db)
    colpali_docs = docs_retrieval_model.search(slide_headings_text, k=3)

    if not colpali_docs:
        raise ValueError("No documents retrieved by Colpali. Try revising the input content.")

    return text_context, colpali_docs

# ==================================================
# 6. Image Processing & Description
# ==================================================
def generate_image_description(img, pipe, processor):
    chat_template = [
        {"role": "user", "content": [
            {"type": "image", "image": img},
            {"type": "text", "text": "Briefly describe the image."}
        ]}
    ]
    prompt = processor.apply_chat_template(chat_template, add_generation_prompt=True)
    outputs = pipe(img, prompt=prompt, generate_kwargs={"max_new_tokens": 200})
    return outputs[0]["generated_text"].split("ASSISTANT:")[-1].strip()

def get_combined_image_context(colpali_docs, all_images, pipe, processor):
    image_contexts = []
    image_path_map = {}
    for idx, result in enumerate(colpali_docs):
        doc_id, page_num = result.doc_id, result.page_num
        image_files = all_images.get(doc_id, [])
        if page_num - 1 < len(image_files):
            img_path = image_files[page_num - 1]
            description = generate_image_description(img_path, pipe, processor)
            image_contexts.append(description)
            image_path_map[idx] = img_path  # Associate slide index with image path

    return " ".join(image_contexts), image_path_map

# ==================================================
# 7. Slide JSON Generation via Claude
# ==================================================
def generate_slides_json(structured_input, text_context, image_context):
    slides_json = call_claude(
        query_or_answers=structured_input,
        context=text_context,
        image_context=image_context,
        mode="slides"
    )
    return json.loads(slides_json)

# ==================================================
# 8. PowerPoint Presentation Creation
# ==================================================
def create_presentation_from_json(json_slides, output_folder, image_path_map=None):
    """
    Generates and saves a PowerPoint presentation based on the slide JSON.
    Supports title and content slides with optional images.

    Args:
        json_slides (list): List of dictionaries describing each slide.
        output_folder (str): Path to folder where the PPTX will be saved.
        image_path_map (dict): Optional. Maps image_index to actual image path.
    
    Returns:
        str: Path to generated presentation.
    """
    prs = Presentation()

    for slide in json_slides:
        is_title = slide.get("is_title_slide") == "yes"
        layout = prs.slide_layouts[0] if is_title else prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)

        # --- Title Slide ---
        if is_title:
            s.shapes.title.text = slide.get("title_text", "")
            if len(s.placeholders) > 1:
                s.placeholders[1].text = slide.get("subtitle_text", "")

        # --- Content Slide ---
        else:
            s.shapes.title.text = slide.get("title_text", "")
            if len(s.placeholders) > 1:
                body = s.placeholders[1].text_frame
                body.clear()
                for bullet in slide.get("text", []):
                    p = body.add_paragraph()
                    p.text = bullet
                    p.level = 0

            # --- Image Insertion ---
            image_index = slide.get("image_index")
            if image_path_map and image_index is not None:
                img_path = image_path_map.get(image_index)
                if img_path and os.path.exists(img_path):
                    # You can adjust placement here
                    s.shapes.add_picture(img_path, Inches(5.5), Inches(1.5), width=Inches(3))

    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, "Generated_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    return output_path


# ==================================================
# 9. DOCX GRANT GENERATION FUNCTION
# ==================================================
def create_universal_grant_docx(
    template_name: str,
    field_ordering: Dict[str, List[str]],
    filled_fields: Dict[str, str],
    title: str = "GRANT PROPOSAL",
    output_filename: str = "grant_proposal.docx"
) -> str:
    if not output_filename.endswith(".docx"):
        output_filename += ".docx"

    doc = Document()
    doc.add_heading(title, level=0)

    roman_numerals = [
        "I", "II", "III", "IV", "V", "VI", "VII",
        "VIII", "IX", "X", "XI", "XII", "XIII", "XIV"
    ]
    ordered_fields = field_ordering.get(template_name, [])

    for idx, field_key in enumerate(ordered_fields):
        content = filled_fields.get(field_key, "")
        if not content:
            continue

        section_num = roman_numerals[idx] if idx < len(roman_numerals) else str(idx + 1)
        section_title = field_key.replace("_", " ").title()

        doc.add_heading(f"{section_num}. {section_title}", level=1)

        if isinstance(content, str):
            doc.add_paragraph(content)
        elif isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    for k, v in item.items():
                        doc.add_paragraph(k.replace("_", " ").title(), style="Heading 2")
                        doc.add_paragraph(str(v))
        elif isinstance(content, dict):
            for k, v in content.items():
                doc.add_paragraph(k.replace("_", " ").title(), style="Heading 2")
                doc.add_paragraph(str(v))

    doc.save(output_filename)
    return os.path.abspath(output_filename)

# ==================================================
# MAIN SLIDE GENERATION FUNCTION
# ==================================================
def generate_slides_from_headings(structured_input):
    all_images, file_names = convert_pdfs_if_needed(DATA_PATH, IMAGES_FOLDER)
    docs_retrieval_model, pipe, processor = initialize_models()
    docs_retrieval_model = index_documents_if_needed(docs_retrieval_model, DATA_PATH, "image_index", force_reindex=True)
    db = create_or_load_vector_db(DATA_PATH, FAISS_DB_PATH, force_rebuild=False)

    purpose = structured_input.get("category", "")
    subtype = structured_input.get("subtype", "")
    answers = structured_input.get("answers", {})

    query_string = f"{purpose}\n{subtype}\n" + "\n".join(answers.values())

    text_context, colpali_docs = retrieve_text_context(query_string, db, docs_retrieval_model)
    image_context, image_path_map = get_combined_image_context(colpali_docs, all_images, pipe, processor)

    slides_data = generate_slides_json(structured_input, text_context, image_context)
    output_path = create_presentation_from_json(slides_data, OUTPUT_FOLDER, image_path_map=image_path_map)
    return output_path


# ==================================================
# MAIN GRANT JSON GENERATION FUNCTION
# ==================================================
def generate_grant_from_inputs(user_prompt_json):
    all_images, _ = convert_pdfs_if_needed(DATA_PATH, IMAGES_FOLDER)
    docs_retrieval_model, pipe, processor = initialize_models()
    docs_retrieval_model = index_documents_if_needed(docs_retrieval_model, DATA_PATH, "image_index", force_reindex=True)
    db = create_or_load_vector_db(DATA_PATH, FAISS_DB_PATH, force_rebuild=False)

    field_values = list(user_prompt_json.get("fields", {}).values())

    if not field_values or all(v == "" for v in field_values):
        raise ValueError("Input fields are empty. Cannot generate query.")

    query_string = "\n".join(
        v if isinstance(v, str) else json.dumps(v, indent=2)
        for v in field_values
    )

    if not query_string.strip():
        raise ValueError("Query string for retrieval is empty.")

    print("[DEBUG] Query string preview:", query_string[:200])

    text_context, colpali_docs = retrieve_text_context(query_string, db, docs_retrieval_model)
    print("Query string:", query_string[:200])
    print("Retrieved Colpali docs:", len(colpali_docs))

    image_context, _ = get_combined_image_context(colpali_docs, all_images, pipe, processor)

    template_name = user_prompt_json.get("template_name", "")
    if not template_name:
        raise ValueError("Missing template name in JSON input.")

    with open("template_fields.json") as f:
        template_fields = json.load(f)

    output_json = call_claude(
        query_or_answers=user_prompt_json["fields"],
        template_name=template_name,
        template_fields=template_fields,
        context=text_context,
        image_context=image_context,
        mode="grant",
        debug=os.getenv("DEBUG", "False") == "True"
    )

    if not output_json:
        raise ValueError("Claude did not return any output. Please retry.")

    try:
        filled_fields = json.loads(output_json)
    except json.JSONDecodeError as e:
        print("[ERROR] JSON decode failed. Raw response:")
        print(output_json)
        raise e

    docx_path = create_universal_grant_docx(
        template_name=template_name,
        field_ordering=template_fields,
        filled_fields=filled_fields,
        title="GRANT PROPOSAL",
        output_filename="Final_Grant_Proposal.docx"
    )

    return docx_path
