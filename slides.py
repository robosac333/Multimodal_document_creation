import os
import json
from pptx import Presentation
from pptx.util import Inches
from text_retrieval import create_vector_db
from image_retrieval import convert_pdfs_to_images, load_existing_image_mappings
from Chatbot.Mistral_7b import retrieve_faiss, retrieve_context
from transformers import pipeline, AutoProcessor
from byaldi import RAGMultiModalModel
from bedrock_handler import call_claude
from create_grant_proposal import create_grant_proposal_from_json
from create_rfp_proposal import create_rfp_proposal_from_json

# ==================================================
# CONFIGURATION & PATHS
# ==================================================
DATA_PATH = "/nfshomes/sjd3333/Retrieval_based_ppt_creation/pdfs"
FAISS_DB_PATH = "/nfshomes/sjd3333/Retrieval_based_ppt_creation/vector_data_base"
IMAGES_FOLDER = "/nfshomes/sjd3333/Retrieval_based_ppt_creation/img_output"
OUTPUT_FOLDER = "/nfshomes/sjd3333/Retrieval_based_ppt_creation/Final_slide"

# ==================================================
# 1. PDF Conversion Caching
# ==================================================
def convert_pdfs_if_needed(data_path, images_folder):
    """
    Converts PDFs to images if necessary. 
    If the images folder does not exist or is empty, it converts the PDFs.
    Otherwise, it loads the existing image mappings.
    """
    if not os.path.exists(images_folder) or not os.listdir(images_folder):
        all_images, file_names = convert_pdfs_to_images(data_path, images_folder)
        print("PDFs converted to images.")
    else:
        # load_existing_image_mappings should reconstruct the mapping from your images folder.
        all_images = load_existing_image_mappings(images_folder)
        file_names = list(all_images.keys())
        print("Using cached image files.")
    return all_images, file_names

# ==================================================
# 2. Document Indexing Caching
# ==================================================
def index_documents_if_needed(model, data_path, index_name, force_reindex=False):
    """
    Indexes documents with the retrieval model if needed.
    Checks for an index file indicator; if it doesn't exist or force_reindex is True,
    indexing is performed.
    """
    index_indicator = os.path.join(data_path, f"{index_name}_index_indicator.txt")
    if force_reindex or not os.path.exists(index_indicator):
        model.index(input_path=data_path, index_name=index_name, overwrite=True)
        # Create an indicator file to signal that indexing has been done.
        with open(index_indicator, 'w') as f:
            f.write("indexed")
        print("Documents indexed.")
    else:
        print("Using cached document index.")
    return model

# ==================================================
# 3. Vector Database Creation/Loading Caching
# ==================================================
def create_or_load_vector_db(data_path, faiss_db_path, force_rebuild=False):
    """
    Creates a new vector database if needed, otherwise loads an existing one.
    """
    if force_rebuild or not os.path.exists(faiss_db_path) or not os.listdir(faiss_db_path):
        os.makedirs(faiss_db_path, exist_ok=True)
        create_vector_db(data_path, faiss_db_path)
        print("Vector database created.")
    else:
        print("Using existing vector database.")
    db = retrieve_faiss(faiss_db_path)
    return db

# ==================================================
# 4. Model and Pipeline Initialization (Without Caching)
# ==================================================
def initialize_models():
    """
    Initializes and returns the document retrieval model, image-to-text pipeline,
    and processor. This does not cache the loaded models.
    """
    docs_retrieval_model = RAGMultiModalModel.from_pretrained("vidore/colpali-v1.2")
    model_id = "llava-hf/llava-1.5-7b-hf"
    pipe = pipeline("image-to-text", model=model_id, device=0)
    processor = AutoProcessor.from_pretrained(model_id)
    return docs_retrieval_model, pipe, processor

# ==================================================
# 5. Text Context Retrieval
# ==================================================
def retrieve_text_context(slide_headings_text, db, docs_retrieval_model):
    """
    Retrieves textual context and document search results based on the slide headings.
    """
    text_context = retrieve_context(slide_headings_text, db)
    colpali_docs = docs_retrieval_model.search(slide_headings_text, k=3)
    return text_context, colpali_docs

# ==================================================
# 6. Image Processing & Description
# ==================================================
def generate_image_description(img, pipe, processor):
    """
    Uses the image-to-text pipeline to generate a description for the given image.
    """
    chat_template = [
        {"role": "user", "content": [
            {"type": "image", "image": img},
            {"type": "text", "text": "Briefly describe the image."}
        ]}
    ]
    prompt = processor.apply_chat_template(chat_template, add_generation_prompt=True)
    outputs = pipe(img, prompt=prompt, generate_kwargs={"max_new_tokens": 200})
    return outputs[0]["generated_text"].split("ASSISTANT:")[-1].strip()

def get_combined_image_context(colpali_docs, all_images, pipe, processor):
    """
    For each document found by the search, retrieves the corresponding image,
    generates a description, and returns a combined string of image contexts.
    """
    image_contexts = []
    for result in colpali_docs:
        doc_id, page_num = result.doc_id, result.page_num
        image_files_for_doc = all_images.get(doc_id, [])
        if page_num - 1 < len(image_files_for_doc):
            img = image_files_for_doc[page_num - 1]
            description = generate_image_description(img, pipe, processor)
            image_contexts.append(description)
    return " ".join(image_contexts)

# ==================================================
# 7. Slide Generation JSON via Claude
# ==================================================
def generate_slides_json(slide_headings_text, text_context, image_context):
    """
    Calls the Claude service to generate a JSON structure for the slides
    based on the slide headings, text context, and image context.
    """
    slides_json = call_claude(
        query=slide_headings_text,
        context=text_context,
        image_context=image_context,
        mode="slides"
    )
    return json.loads(slides_json)

# ==================================================
# 8. Report Generation JSON via Claude
# ==================================================
def generate_report_json(slide_headings_text, text_context, image_context):
    """
    Calls the Claude service to generate a JSON structure for the slides
    based on the slide headings, text context, and image context.
    """
    slides_json = call_claude(
        query=slide_headings_text,
        context=text_context,
        image_context=image_context,
        mode="report"
    )
    return json.loads(slides_json)

# ==================================================
# 9. PowerPoint Presentation Creation
# ==================================================
def create_presentation_from_json(json_slides, output_folder):
    """
    Generates and saves a PowerPoint presentation based on the slide JSON.
    Returns the path to the generated presentation.
    """
    prs = Presentation()
    for slide in json_slides:
        if slide.get("is_title_slide") == "yes":
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title_text"]
            s.placeholders[1].text = slide["subtitle_text"]
        else:
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title_text"]
            body = s.placeholders[1].text_frame
            body.clear()
            for bullet in slide["text"]:
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0
    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, "Generated_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    return output_path


def load_vector_database():
    # Step 1: Convert PDFs to images (with caching)
    all_images, file_names = convert_pdfs_if_needed(DATA_PATH, IMAGES_FOLDER)

    # Step 2: Initialize models and pipelines (models loaded freshly)
    docs_retrieval_model, pipe, processor = initialize_models()

    # Step 3: Index the documents if needed (with caching)
    docs_retrieval_model = index_documents_if_needed(docs_retrieval_model, DATA_PATH, "image_index", force_reindex=True)

    # Step 4: Create or load the vector database (with caching)
    db = create_or_load_vector_db(DATA_PATH, FAISS_DB_PATH, force_rebuild=True)

    return db, docs_retrieval_model, all_images, pipe, processor

# ==================================================
# MAIN SLIDE GENERATION FUNCTION
# ==================================================
def generate_slides_from_headings(db, docs_retrieval_model, all_images, pipe, processor, slide_headings_text):
    """
    Main function that, given slide headings as input, processes the PDFs,
    retrieves contexts, generates slide JSON, and creates a PowerPoint presentation.
    Returns the path to the generated presentation for downloading.
    """
    # Step 5: Retrieve text context and search for relevant documents
    text_context, colpali_docs = retrieve_text_context(slide_headings_text, db, docs_retrieval_model)

    # Step 6: Process images to generate a combined image context string
    image_context = get_combined_image_context(colpali_docs, all_images, pipe, processor)

    # Step 7: Generate slide JSON using Claude
    slides_data = generate_slides_json(slide_headings_text, text_context, image_context)

    # Step 8: Create the PowerPoint presentation from the JSON data
    output_path = create_presentation_from_json(slides_data, OUTPUT_FOLDER)
    return output_path

# ==================================================
# RESEARCH GRANT PROPOSAL FUNCTION
# ==================================================
def generate_grant_proposal_from_headings(db, docs_retrieval_model, all_images, pipe, processor, report_json, template_type):
    """
    Main function that, given slide headings as input, processes the PDFs,
    retrieves contexts, generates slide JSON, and creates a PowerPoint presentation.
    Returns the path to the generated presentation for downloading.
    """
    # Step 5: Retrieve text context and search for relevant documents
    text_context, colpali_docs = retrieve_text_context(report_json, db, docs_retrieval_model)

    # Step 6: Process images to generate a combined image context string
    image_context = get_combined_image_context(colpali_docs, all_images, pipe, processor)

    # Step 7: Generate slide JSON using Claude
    report_json = generate_report_json(report_json, text_context, image_context)

    # Step 8: Create the PowerPoint presentation from the JSON data
    if template_type == "proposal":
        output_path = create_grant_proposal_from_json(report_json)
    elif template_type == "rfp":
        output_path = create_rfp_proposal_from_json(report_json)

    return output_path

# # ==================================================
# # OPTIONAL: Standalone script testing
# # ==================================================
# if __name__ == "__main__":
#     # For testing purposes, use a sample slide headings string.
#     sample_headings = "Introduction\nMarket Analysis\nFinancial Overview\nConclusion"
#     ppt_path = generate_slides_from_headings(sample_headings)
#     print(f"Generated presentation available at: {ppt_path}")
